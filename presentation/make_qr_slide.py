# -*- coding: utf-8 -*-
"""발표 첫 화면 PPT 1장 — QR로 웹사이트 접속."""
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent
URL = "https://inspector11.vercel.app/"
QR_PATH = ROOT / "qr_inspector11.png"
PPTX_PATH = ROOT / "11조_함께하조_QR접속.pptx"
EMBLEM = REPO / "assets_emblem_only.svg"
EMBLEM_PNG = REPO / "assets_emblem_only.png"

NAVY = RGBColor(0x02, 0x32, 0x5C)
NAVY2 = RGBColor(0x0A, 0x4A, 0x82)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
GOV_RED = RGBColor(0xE6, 0x00, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xDC, 0xE8, 0xF6)


def make_qr():
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=14,
        border=2,
    )
    qr.add_data(URL)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#02325C", back_color="white").convert("RGB")
    img.save(QR_PATH)
    print("QR:", QR_PATH)


def set_run(run, size, bold=False, color=WHITE, font="맑은 고딕"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    try:
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            from pptx.oxml import parse_xml
            ea = parse_xml(f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>')
            rPr.append(ea)
        else:
            ea.set("typeface", font)
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, text, size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=bold, color=color)
    return box


def main():
    make_qr()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # 상단 정부 레드 라인
    red = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    red.fill.solid()
    red.fill.fore_color.rgb = GOV_RED
    red.line.fill.background()

    # 골드 악센트 바
    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.35), Inches(0.12), Inches(1.1)
    )
    gold.fill.solid()
    gold.fill.fore_color.rgb = GOLD
    gold.line.fill.background()

    # 엠블럼
    emblem = EMBLEM_PNG if EMBLEM_PNG.exists() else None
    if emblem:
        slide.shapes.add_picture(str(emblem), Inches(0.9), Inches(0.35), height=Inches(0.55))

    add_textbox(
        slide, Inches(1.6), Inches(0.38), Inches(6), Inches(0.5),
        "노사누리 · 근로감독행정시스템", 16, False, SOFT, PP_ALIGN.LEFT
    )

    add_textbox(
        slide, Inches(1.15), Inches(1.25), Inches(7.2), Inches(0.7),
        "11조  함께하조", 40, True, WHITE, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(1.15), Inches(1.95), Inches(7.2), Inches(0.55),
        "노동감독관, 일의 의미", 22, False, GOLD, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(1.15), Inches(2.7), Inches(7.0), Inches(1.2),
        "휴대폰 카메라로 QR을 스캔하면\n지금 바로 조 소개 웹페이지에 접속할 수 있습니다.",
        18, False, SOFT, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(1.15), Inches(4.1), Inches(7.0), Inches(0.45),
        "접속 후 시연 영상을 함께 봐 주세요", 16, False, WHITE, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(1.15), Inches(6.55), Inches(7.5), Inches(0.4),
        URL, 14, False, SOFT, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(1.15), Inches(6.95), Inches(7.5), Inches(0.35),
        "신규 노동감독관 수사학교  ·  조별 발표", 12, False, RGBColor(0x9F, 0xB8, 0xD4), PP_ALIGN.LEFT
    )

    # QR 카드 (흰 배경)
    card_l, card_t = Inches(8.55), Inches(1.35)
    card_w, card_h = Inches(3.9), Inches(4.85)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_l, card_t, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    try:
        card.adjustments[0] = 0.08
    except Exception:
        pass

    qr_size = Inches(3.15)
    qr_left = card_l + (card_w - qr_size) / 2
    slide.shapes.add_picture(str(QR_PATH), qr_left, card_t + Inches(0.35), width=qr_size)

    add_textbox(
        slide, card_l, card_t + Inches(3.65), card_w, Inches(0.4),
        "QR 스캔", 18, True, NAVY, PP_ALIGN.CENTER
    )
    add_textbox(
        slide, card_l + Inches(0.15), card_t + Inches(4.1), card_w - Inches(0.3), Inches(0.55),
        "inspector11.vercel.app", 12, False, NAVY2, PP_ALIGN.CENTER
    )

    prs.save(PPTX_PATH)
    print("PPTX:", PPTX_PATH)


if __name__ == "__main__":
    main()
